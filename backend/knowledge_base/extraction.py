"""Text extraction from uploaded documents."""

import io

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}

# PDFs with less than this many characters of pdfminer output are treated as scanned
PDF_OCR_FALLBACK_THRESHOLD = 200
PDF_OCR_MAX_PAGES = 20


def extract_text(file_obj, ext: str) -> tuple[str | None, str]:
    """
    Returns (text, status) where status is 'indexed' | 'reference' | 'error'.
    text is None for binary/reference files.
    """
    ext = ext.lower()

    try:
        if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html"):
            raw = file_obj.read()
            return raw.decode("utf-8", errors="replace"), "indexed"

        if ext in (".xlsx", ".xls", ".xlsm"):
            return _extract_excel(file_obj), "indexed"

        if ext == ".pdf":
            return _extract_pdf(file_obj), "indexed"

        if ext in (".docx",):
            return _extract_docx(file_obj), "indexed"

        if ext == ".pptx":
            return _extract_pptx(file_obj), "indexed"

        if ext in IMAGE_EXTS:
            return _extract_image(file_obj), "indexed"

        # Binary formats we can't parse
        return None, "reference"

    except Exception as exc:
        return f"[Extractiefout: {exc}]", "error"


def _ocr(image_bytes: bytes) -> str:
    import pytesseract
    from PIL import Image

    image = Image.open(io.BytesIO(image_bytes))
    return pytesseract.image_to_string(image, lang="nld+eng")


def _extract_image(file_obj) -> str:
    return _ocr(file_obj.read())


def _extract_pdf(file_obj) -> str:
    from pdfminer.high_level import extract_text_to_fp
    from pdfminer.layout import LAParams

    file_bytes = file_obj.read()

    # Try pdfminer first (fast, free, works for text-based PDFs)
    out = io.StringIO()
    extract_text_to_fp(
        io.BytesIO(file_bytes),
        out,
        laparams=LAParams(),
        output_type="text",
        codec="utf-8",
    )
    text = out.getvalue().strip()

    if len(text) >= PDF_OCR_FALLBACK_THRESHOLD:
        return text

    # Fallback: render pages as images and OCR via Claude (for scanned PDFs)
    return _extract_pdf_ocr(file_bytes)


def _extract_pdf_ocr(file_bytes: bytes) -> str:
    import fitz  # PyMuPDF

    pdf = fitz.open(stream=file_bytes, filetype="pdf")
    page_count = min(len(pdf), PDF_OCR_MAX_PAGES)
    all_text = []

    for page_num in range(page_count):
        page = pdf[page_num]
        mat = fitz.Matrix(2, 2)  # 2× zoom for better OCR quality
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")
        page_text = _ocr(img_bytes)
        if page_text.strip():
            all_text.append(f"[Pagina {page_num + 1}]\n{page_text}")

    pdf.close()
    return "\n\n".join(all_text)


def _extract_docx(file_obj) -> str:
    from docx import Document

    doc = Document(file_obj)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            paragraphs.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(paragraphs)


def _extract_pptx(file_obj) -> str:
    from pptx import Presentation

    prs = Presentation(file_obj)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Dia {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _extract_excel(file_obj) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        all_rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                all_rows.append(cells)

        if not all_rows:
            continue

        lines.append(f"[Tabblad: {sheet.title} — {len(all_rows)} rijen]")

        # Detect header row: first row where most cells are non-numeric text
        first = all_rows[0]
        non_empty = [c.strip() for c in first if c.strip()]
        numeric = sum(1 for c in non_empty if _looks_numeric(c))
        has_headers = non_empty and numeric < len(non_empty) / 2

        if has_headers and len(all_rows) > 1:
            headers = [c.strip() for c in first]
            lines.append(f"[Kolommen: {' | '.join(h for h in headers if h)}]")
            for row_cells in all_rows[1:]:
                if any(c.strip() for c in row_cells):
                    pairs = [
                        f"{h}: {v}"
                        for h, v in zip(headers, row_cells)
                        if h and v.strip()
                    ]
                    if pairs:
                        lines.append(", ".join(pairs))
        else:
            for row_cells in all_rows:
                lines.append("\t".join(row_cells))

    return "\n".join(lines)


def _looks_numeric(s: str) -> bool:
    s = s.strip().replace(",", ".").replace("%", "").replace("€", "").replace("-", "")
    try:
        float(s)
        return True
    except ValueError:
        return False
